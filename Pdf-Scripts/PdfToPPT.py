import fitz  
import sys
import os
from pptx import Presentation
from pptx.util import Inches

def pdf_to_ppt(pdf_path, out_ppt="output.pptx", zoom=2):
    doc = fitz.open(pdf_path)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  
    mat = fitz.Matrix(zoom, zoom)

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(matrix=mat)
        img_path = f"page_{page_num+1:03d}.jpg"
        pix.save(img_path)

        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            img_path,
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )
        os.remove(img_path)  

    doc.close()
    prs.save(out_ppt)
    return out_ppt


if __name__ == "__main__":
    pdf_file = sys.argv[1]
    out_ppt = sys.argv[2]
    pdf_to_ppt(pdf_file, out_ppt)
